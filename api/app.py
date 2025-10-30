from fastapi import FastAPI, Form, Request
from fastapi.responses import StreamingResponse, HTMLResponse
from fastapi.templating import Jinja2Templates
from pptx import Presentation
from io import BytesIO

app = FastAPI()
templates = Jinja2Templates(directory="../dtemplates")  # مجلد HTML

# =========================
# واجهة المستخدم (HTML)
# =========================
@app.get("/", response_class=HTMLResponse)
def home(request: Request):
    return templates.TemplateResponse("index.html", {"request": request})


# =========================
# توليد ملف PowerPoint
# =========================
@app.post("/generate")
def generate_pptx(
    name: str = Form(...),
    nationality: str = Form(...),
    birthdate: str = Form(...),
    city: str = Form(...),
    id_number: str = Form(...)
):
    # فتح قالب PowerPoint
    prs = Presentation("../template.pptx")
    
    # قاموس الكلمات المفتاحية والقيم المدخلة
    replacements = {
        "{name}": name,
        "{nationality}": nationality,
        "{birthdate}": birthdate,
        "{city}": city,
        "{id_number}": id_number
    }

    # المرور على كل الشرائح والأشكال والنصوص
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in replacements.items():
                            if key in run.text:
                                run.text = run.text.replace(key, value)

    # حفظ الملف في الذاكرة
    output = BytesIO()
    prs.save(output)
    output.seek(0)

    # إعادة الملف للتحميل مباشرة
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=generated.pptx"}
    )